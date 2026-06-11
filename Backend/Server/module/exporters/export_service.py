"""
Export Service for MODULE
=========================

Exports lessons and assignments to PDF, DOC, and PPT formats.
"""

import io
import base64
import time
from typing import Optional
from datetime import datetime

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    Image as RLImage, PageBreak, ListFlowable, ListItem
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN

from ..models.schemas import (
    Lesson, Assignment, Slide, DifficultyLevel, QuestionType,
    MCQQuestion, ShortAnswerQuestion, LongAnswerQuestion
)


class ExportService:
    """Exports lessons and assignments to PDF, DOC, PPT formats."""
    
    def __init__(self):
        """Initialize export service with default styles."""
        self._setup_pdf_styles()
    
    def _setup_pdf_styles(self):
        """Setup PDF styles for consistent formatting."""
        self.styles = getSampleStyleSheet()
        
        # Custom title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor=colors.HexColor('#1a365d')
        ))
        
        # Slide title style
        self.styles.add(ParagraphStyle(
            name='SlideTitle',
            parent=self.styles['Heading2'],
            fontSize=16,
            spaceBefore=20,
            spaceAfter=10,
            textColor=colors.HexColor('#2c5282')
        ))
        
        # Content style
        self.styles.add(ParagraphStyle(
            name='Content',
            parent=self.styles['Normal'],
            fontSize=11,
            spaceAfter=8,
            alignment=TA_JUSTIFY
        ))
        
        # Custom bullet style (renamed to avoid conflict)
        self.styles.add(ParagraphStyle(
            name='CustomBullet',
            parent=self.styles['Normal'],
            fontSize=10,
            leftIndent=20,
            spaceAfter=4
        ))
        
        # Key term style
        self.styles.add(ParagraphStyle(
            name='KeyTerm',
            parent=self.styles['Normal'],
            fontSize=10,
            textColor=colors.HexColor('#2d3748'),
            backColor=colors.HexColor('#edf2f7'),
            leftIndent=10,
            rightIndent=10,
            spaceBefore=2,
            spaceAfter=2
        ))

    async def export_lesson_pdf(
        self,
        lesson: Lesson,
        include_diagrams: bool = True
    ) -> bytes:
        """
        Export lesson as PDF with embedded diagrams.
        
        Args:
            lesson: The lesson to export
            include_diagrams: Whether to include diagram images
            
        Returns:
            PDF file as bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch
        )
        
        story = []
        
        # Title page
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(lesson.topic, self.styles['CustomTitle']))
        story.append(Spacer(1, 0.5*inch))
        
        # Metadata
        meta_data = [
            ['Class:', lesson.class_name],
            ['Subject:', lesson.subject],
            ['Created:', lesson.created_at.strftime('%B %d, %Y')]
        ]
        meta_table = Table(meta_data, colWidths=[1.5*inch, 4*inch])
        meta_table.setStyle(TableStyle([
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))
        story.append(meta_table)
        story.append(PageBreak())
        
        # Each slide
        for slide in lesson.slides:
            story.extend(self._create_pdf_slide(slide, include_diagrams))
            story.append(PageBreak())
        
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _create_pdf_slide(
        self,
        slide: Slide,
        include_diagrams: bool
    ) -> list:
        """Create PDF elements for a single slide."""
        elements = []
        
        # Slide header
        slide_header = f"Slide {slide.slide_number}: {slide.title}"
        elements.append(Paragraph(slide_header, self.styles['SlideTitle']))
        elements.append(Spacer(1, 0.2*inch))
        
        # Explanation
        elements.append(Paragraph("<b>Explanation:</b>", self.styles['Content']))
        elements.append(Paragraph(slide.explanation, self.styles['Content']))
        elements.append(Spacer(1, 0.15*inch))
        
        # Bullet points
        if slide.bullet_points:
            elements.append(Paragraph("<b>Key Points:</b>", self.styles['Content']))
            bullet_items = []
            for point in slide.bullet_points:
                bullet_items.append(ListItem(
                    Paragraph(point, self.styles['CustomBullet']),
                    leftIndent=20
                ))
            elements.append(ListFlowable(bullet_items, bulletType='bullet'))
            elements.append(Spacer(1, 0.15*inch))
        
        # Key terms
        if slide.key_terms:
            elements.append(Paragraph("<b>Key Terms:</b>", self.styles['Content']))
            for term in slide.key_terms:
                elements.append(Paragraph(f"• {term}", self.styles['KeyTerm']))
            elements.append(Spacer(1, 0.15*inch))
        
        # Examples
        if slide.examples:
            elements.append(Paragraph("<b>Examples:</b>", self.styles['Content']))
            for i, example in enumerate(slide.examples, 1):
                elements.append(Paragraph(f"{i}. {example}", self.styles['Content']))
            elements.append(Spacer(1, 0.15*inch))
        
        # Diagram
        if include_diagrams and slide.diagram_url:
            elements.append(Paragraph("<b>Diagram:</b>", self.styles['Content']))
            try:
                # Handle base64 encoded images
                if slide.diagram_url.startswith('data:image'):
                    img_data = slide.diagram_url.split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    img_buffer = io.BytesIO(img_bytes)
                    img = RLImage(img_buffer, width=4*inch, height=3*inch)
                    elements.append(img)
                else:
                    # URL-based image (would need to fetch)
                    elements.append(Paragraph(
                        f"[Diagram: {slide.diagram_prompt}]",
                        self.styles['Content']
                    ))
            except Exception:
                elements.append(Paragraph(
                    f"[Diagram placeholder: {slide.diagram_prompt}]",
                    self.styles['Content']
                ))
        elif slide.diagram_prompt:
            elements.append(Paragraph(
                f"<b>Diagram Prompt:</b> {slide.diagram_prompt}",
                self.styles['Content']
            ))
        
        return elements

    async def export_lesson_doc(
        self,
        lesson: Lesson,
        include_diagrams: bool = True
    ) -> bytes:
        """
        Export lesson as Word document with embedded diagrams.
        
        Args:
            lesson: The lesson to export
            include_diagrams: Whether to include diagram images
            
        Returns:
            DOCX file as bytes
        """
        doc = Document()
        
        # Title
        title = doc.add_heading(lesson.topic, level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata
        meta_para = doc.add_paragraph()
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        meta_para.add_run(f"Class: {lesson.class_name}").bold = True
        meta_para.add_run(f"  |  Subject: {lesson.subject}")
        meta_para.add_run(f"  |  Created: {lesson.created_at.strftime('%B %d, %Y')}")
        
        doc.add_paragraph()  # Spacer
        
        # Each slide
        for slide in lesson.slides:
            self._add_doc_slide(doc, slide, include_diagrams)
            doc.add_page_break()
        
        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _add_doc_slide(
        self,
        doc: Document,
        slide: Slide,
        include_diagrams: bool
    ):
        """Add a slide to the Word document."""
        # Slide header
        heading = doc.add_heading(f"Slide {slide.slide_number}: {slide.title}", level=1)
        
        # Explanation
        doc.add_heading("Explanation", level=2)
        doc.add_paragraph(slide.explanation)
        
        # Bullet points
        if slide.bullet_points:
            doc.add_heading("Key Points", level=2)
            for point in slide.bullet_points:
                para = doc.add_paragraph(point, style='List Bullet')
        
        # Key terms
        if slide.key_terms:
            doc.add_heading("Key Terms", level=2)
            for term in slide.key_terms:
                para = doc.add_paragraph()
                para.add_run("• ").bold = True
                para.add_run(term)
        
        # Examples
        if slide.examples:
            doc.add_heading("Examples", level=2)
            for i, example in enumerate(slide.examples, 1):
                doc.add_paragraph(f"{i}. {example}")
        
        # Diagram
        if include_diagrams and slide.diagram_url:
            doc.add_heading("Diagram", level=2)
            try:
                if slide.diagram_url.startswith('data:image'):
                    img_data = slide.diagram_url.split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    img_buffer = io.BytesIO(img_bytes)
                    doc.add_picture(img_buffer, width=Inches(5))
                else:
                    doc.add_paragraph(f"[Diagram: {slide.diagram_prompt}]")
            except Exception:
                doc.add_paragraph(f"[Diagram placeholder: {slide.diagram_prompt}]")
        elif slide.diagram_prompt:
            doc.add_heading("Diagram Prompt", level=2)
            doc.add_paragraph(slide.diagram_prompt)

    async def export_lesson_ppt(
        self,
        lesson: Lesson,
        include_diagrams: bool = True
    ) -> bytes:
        """
        Export lesson as PowerPoint presentation.
        
        Args:
            lesson: The lesson to export
            include_diagrams: Whether to include diagram images
            
        Returns:
            PPTX file as bytes
        """
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = lesson.topic
        subtitle.text = f"Class: {lesson.class_name} | Subject: {lesson.subject}\n{lesson.created_at.strftime('%B %d, %Y')}"
        
        # Content slides
        for lesson_slide in lesson.slides:
            self._add_ppt_slide(prs, lesson_slide, include_diagrams)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _add_ppt_slide(
        self,
        prs: Presentation,
        slide: Slide,
        include_diagrams: bool
    ):
        """Add a slide to the PowerPoint presentation."""
        # Use content layout
        content_layout = prs.slide_layouts[1]  # Title and Content
        ppt_slide = prs.slides.add_slide(content_layout)
        
        # Title
        title_shape = ppt_slide.shapes.title
        title_shape.text = f"Slide {slide.slide_number}: {slide.title}"
        
        # Content placeholder
        content_shape = ppt_slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()
        
        # Add explanation (truncated for slide)
        p = tf.paragraphs[0]
        p.text = slide.explanation[:200] + "..." if len(slide.explanation) > 200 else slide.explanation
        p.font.size = PptxPt(14)
        
        # Add bullet points
        for point in slide.bullet_points[:5]:  # Limit to 5 for readability
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
            p.font.size = PptxPt(12)
        
        # Add key terms (if space allows)
        if slide.key_terms:
            p = tf.add_paragraph()
            p.text = "Key Terms: " + ", ".join(slide.key_terms[:3])
            p.level = 0
            p.font.size = PptxPt(11)
            p.font.italic = True
        
        # Add diagram if available
        if include_diagrams and slide.diagram_url:
            try:
                if slide.diagram_url.startswith('data:image'):
                    img_data = slide.diagram_url.split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    img_buffer = io.BytesIO(img_bytes)
                    # Add image to right side of slide
                    ppt_slide.shapes.add_picture(
                        img_buffer,
                        PptxInches(8),
                        PptxInches(1.5),
                        width=PptxInches(4.5)
                    )
            except Exception:
                pass  # Skip diagram if it fails

    async def export_assignment_pdf(
        self,
        assignment: Assignment,
        include_answers: bool = False
    ) -> bytes:
        """
        Export assignment as PDF, optionally with answer key.
        
        Args:
            assignment: The assignment to export
            include_answers: Whether to include answer key
            
        Returns:
            PDF file as bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.75*inch,
            bottomMargin=0.75*inch
        )
        
        story = []
        
        # Title
        story.append(Paragraph(f"Assignment: {assignment.topic}", self.styles['CustomTitle']))
        story.append(Spacer(1, 0.3*inch))
        
        # Metadata
        meta_text = f"Class: {assignment.class_name} | Subject: {assignment.subject} | Total Marks: {assignment.total_marks}"
        story.append(Paragraph(meta_text, self.styles['Content']))
        story.append(Spacer(1, 0.5*inch))
        
        # Group questions by difficulty
        questions_by_diff = assignment.questions_by_difficulty()
        question_num = 1
        
        for difficulty in DifficultyLevel:
            questions = questions_by_diff.get(difficulty, [])
            if not questions:
                continue
            
            # Section header
            story.append(Paragraph(
                f"<b>Section: {difficulty.value.upper()} Questions</b>",
                self.styles['SlideTitle']
            ))
            story.append(Spacer(1, 0.2*inch))
            
            for question in questions:
                story.extend(self._create_pdf_question(
                    question, question_num, include_answers
                ))
                question_num += 1
            
            story.append(Spacer(1, 0.3*inch))
        
        # Answer key section (if requested)
        if include_answers:
            story.append(PageBreak())
            story.append(Paragraph("ANSWER KEY", self.styles['CustomTitle']))
            story.append(Spacer(1, 0.3*inch))
            
            question_num = 1
            for difficulty in DifficultyLevel:
                questions = questions_by_diff.get(difficulty, [])
                for question in questions:
                    story.extend(self._create_pdf_answer(question, question_num))
                    question_num += 1
        
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _create_pdf_question(
        self,
        question,
        question_num: int,
        include_answers: bool
    ) -> list:
        """Create PDF elements for a single question."""
        elements = []
        
        # Question text
        q_text = f"<b>Q{question_num}.</b> {question.question_text} [{question.marks} marks]"
        elements.append(Paragraph(q_text, self.styles['Content']))
        
        # MCQ options
        if question.question_type == QuestionType.MCQ:
            options = ['a', 'b', 'c', 'd']
            for i, opt in enumerate(question.options):
                opt_text = f"    ({options[i]}) {opt.option_text}"
                elements.append(Paragraph(opt_text, self.styles['CustomBullet']))
        
        elements.append(Spacer(1, 0.15*inch))
        return elements
    
    def _create_pdf_answer(self, question, question_num: int) -> list:
        """Create PDF elements for an answer."""
        elements = []
        
        if question.question_type == QuestionType.MCQ:
            correct_opt = None
            options = ['a', 'b', 'c', 'd']
            for i, opt in enumerate(question.options):
                if opt.is_correct:
                    correct_opt = options[i]
                    break
            answer_text = f"<b>Q{question_num}:</b> ({correct_opt})"
        elif question.question_type == QuestionType.SHORT_ANSWER:
            answer_text = f"<b>Q{question_num}:</b> {question.expected_answer}"
        else:  # Long answer
            answer_text = f"<b>Q{question_num}:</b> {question.expected_answer}"
            elements.append(Paragraph(answer_text, self.styles['Content']))
            elements.append(Paragraph("<b>Marking Scheme:</b>", self.styles['Content']))
            for point in question.marking_scheme:
                elements.append(Paragraph(f"• {point}", self.styles['CustomBullet']))
            elements.append(Spacer(1, 0.1*inch))
            return elements
        
        elements.append(Paragraph(answer_text, self.styles['Content']))
        elements.append(Spacer(1, 0.1*inch))
        return elements

    async def export_assignment_doc(
        self,
        assignment: Assignment,
        include_answers: bool = False
    ) -> bytes:
        """
        Export assignment as Word document.
        
        Args:
            assignment: The assignment to export
            include_answers: Whether to include answer key
            
        Returns:
            DOCX file as bytes
        """
        doc = Document()
        
        # Title
        title = doc.add_heading(f"Assignment: {assignment.topic}", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadata
        meta_para = doc.add_paragraph()
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        meta_para.add_run(f"Class: {assignment.class_name}").bold = True
        meta_para.add_run(f"  |  Subject: {assignment.subject}")
        meta_para.add_run(f"  |  Total Marks: {assignment.total_marks}")
        
        doc.add_paragraph()  # Spacer
        
        # Group questions by difficulty
        questions_by_diff = assignment.questions_by_difficulty()
        question_num = 1
        
        for difficulty in DifficultyLevel:
            questions = questions_by_diff.get(difficulty, [])
            if not questions:
                continue
            
            # Section header
            doc.add_heading(f"Section: {difficulty.value.upper()} Questions", level=1)
            
            for question in questions:
                self._add_doc_question(doc, question, question_num)
                question_num += 1
        
        # Answer key section (if requested)
        if include_answers:
            doc.add_page_break()
            answer_title = doc.add_heading("ANSWER KEY", level=0)
            answer_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            question_num = 1
            for difficulty in DifficultyLevel:
                questions = questions_by_diff.get(difficulty, [])
                for question in questions:
                    self._add_doc_answer(doc, question, question_num)
                    question_num += 1
        
        # Save to bytes
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _add_doc_question(self, doc: Document, question, question_num: int):
        """Add a question to the Word document."""
        # Question text
        para = doc.add_paragraph()
        para.add_run(f"Q{question_num}. ").bold = True
        para.add_run(question.question_text)
        para.add_run(f" [{question.marks} marks]").italic = True
        
        # MCQ options
        if question.question_type == QuestionType.MCQ:
            options = ['a', 'b', 'c', 'd']
            for i, opt in enumerate(question.options):
                doc.add_paragraph(f"    ({options[i]}) {opt.option_text}")
        
        doc.add_paragraph()  # Spacer
    
    def _add_doc_answer(self, doc: Document, question, question_num: int):
        """Add an answer to the Word document."""
        para = doc.add_paragraph()
        para.add_run(f"Q{question_num}: ").bold = True
        
        if question.question_type == QuestionType.MCQ:
            options = ['a', 'b', 'c', 'd']
            for i, opt in enumerate(question.options):
                if opt.is_correct:
                    para.add_run(f"({options[i]})")
                    break
        elif question.question_type == QuestionType.SHORT_ANSWER:
            para.add_run(question.expected_answer)
        else:  # Long answer
            para.add_run(question.expected_answer)
            doc.add_paragraph()
            marking_para = doc.add_paragraph()
            marking_para.add_run("Marking Scheme:").bold = True
            for point in question.marking_scheme:
                doc.add_paragraph(f"• {point}", style='List Bullet')
